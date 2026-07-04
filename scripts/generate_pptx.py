"""
Generador de la presentacion PPTX del TFM SauronSheet.
Autor del proyecto: Gonzalo Cantarero Galvez.

Construye un deck completo y bien estructurado a partir de la informacion
del README, cubriendo el funcionamiento de cada pantalla (tab) de la app,
con el sistema de diseno Olive y el logo del proyecto.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Design system (Olive)
# ----------------------------------------------------------------------------
PRIMARY      = RGBColor(0x55, 0x6B, 0x2F)   # Olive Green
PRIMARY_DARK = RGBColor(0x3D, 0x4D, 0x22)
PRIMARY_LT   = RGBColor(0x6B, 0x7B, 0x3A)
ACCENT       = RGBColor(0xD4, 0xA0, 0x17)   # Gold
CANVAS       = RGBColor(0xF8, 0xF9, 0xFA)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK    = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_MUTED   = RGBColor(0x6B, 0x72, 0x80)
INCOME       = RGBColor(0x2E, 0x7D, 0x32)
EXPENSE      = RGBColor(0xC6, 0x28, 0x28)
CARD_BORDER  = RGBColor(0xE5, 0xE0, 0xD5)
CARD_BG      = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
FONT_LIGHT = "Segoe UI Light"

LOGO = r"resources/sauron-sheet-logo-clean.png"
OUT  = r"resources/SauronSheet_TFM_Presentation.pptx"

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def _no_line(shape):
    shape.line.fill.background()

def rect(slide, x, y, w, h, color, line=None, shadow=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    if shadow:
        sp.shadow.inherit = True
    return sp

def rrect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.adjustments[0] = 0.06
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb, tf

def set_run(r, text, size, color, bold=False, italic=False, font=FONT):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font

def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         space_after=6, space_before=0, level=0, font=FONT, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    r = p.add_run()
    set_run(r, text, size, color, bold, italic, font)
    return p

def bg(slide, color=CANVAS):
    rect(slide, 0, 0, SW, SH, color)

def header(slide, title, section=None):
    """Olive header bar with title (and optional section number chip)."""
    bar = rect(slide, 0, 0, SW, Inches(1.05), PRIMARY)
    # accent underline
    rect(slide, 0, Inches(1.05), SW, Inches(0.06), ACCENT)
    tb, tf = textbox(slide, Inches(0.6), Inches(0.02), Inches(11.5), Inches(1.0),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 26, WHITE, bold=True, first=True)
    if section:
        chip = rrect(slide, Inches(12.1), Inches(0.32), Inches(0.7), Inches(0.42), ACCENT)
        ctf = chip.text_frame
        ctf.word_wrap = False
        p = ctf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, str(section), 18, PRIMARY_DARK, bold=True)

def footer(slide, page):
    tb, tf = textbox(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.35))
    para(tf, "SauronSheet · TFM · Gonzalo Cantarero Galvez", 9, TEXT_MUTED, first=True)
    tb2, tf2 = textbox(slide, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35))
    para(tf2, str(page), 9, TEXT_MUTED, align=PP_ALIGN.RIGHT, first=True)

def new_slide():
    return prs.slides.add_slide(BLANK)

# card with title + bullet list
def bullet_card(slide, x, y, w, h, title, items, title_color=PRIMARY,
                icon=None, body_size=13):
    card = rrect(slide, x, y, w, h, CARD_BG, line=CARD_BORDER)
    tb, tf = textbox(slide, x + Inches(0.22), y + Inches(0.14),
                     w - Inches(0.44), h - Inches(0.28))
    head = (icon + "  " if icon else "") + title
    para(tf, head, 15, title_color, bold=True, first=True, space_after=6)
    for it in items:
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        p = tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(3)
        prefix = "   – " if lvl else "•  "
        r = p.add_run()
        set_run(r, prefix + txt, body_size if lvl == 0 else body_size - 1,
                TEXT_DARK if lvl == 0 else TEXT_MUTED)
    return card

def bullets_block(slide, x, y, w, h, items, size=15, color=TEXT_DARK, gap=7):
    tb, tf = textbox(slide, x, y, w, h)
    first = True
    for it in items:
        lvl = 0
        marker = "•  "
        if isinstance(it, tuple):
            txt, lvl = it
            marker = "•  " if lvl == 0 else "   – "
        else:
            txt = it
        p = para(tf, marker + txt, size if lvl == 0 else size - 1,
                 color if lvl == 0 else TEXT_MUTED, first=first, space_after=gap)
        first = False
    return tf

# ----------------------------------------------------------------------------
# SLIDE 1 — Cover
# ----------------------------------------------------------------------------
s = new_slide()
bg(s, PRIMARY_DARK)
# side accent band
rect(s, 0, 0, Inches(0.35), SH, ACCENT)
# logo (landscape ~1.45:1, transparent)
s.shapes.add_picture(LOGO, Inches(0.55), Inches(2.45), height=Inches(2.3))
# title block
tb, tf = textbox(s, Inches(4.1), Inches(1.5), Inches(8.6), Inches(3.2),
                 anchor=MSO_ANCHOR.MIDDLE)
para(tf, "SauronSheet", 54, WHITE, bold=True, first=True, space_after=2)
para(tf, "Control Financiero Personal", 24, ACCENT, bold=True, space_after=14)
para(tf, "Aplicación web de gestión de finanzas personales con importación "
         "de extractos bancarios, análisis avanzado y presupuestos inteligentes.",
     15, RGBColor(0xE8, 0xE8, 0xDD), space_after=0)
# separator line
rect(s, Inches(4.15), Inches(5.05), Inches(8.2), Pt(2), ACCENT)
tb2, tf2 = textbox(s, Inches(4.1), Inches(5.25), Inches(8.6), Inches(1.7))
para(tf2, "Trabajo de Fin de Máster", 15, WHITE, bold=True, first=True, space_after=2)
para(tf2, "Máster de Desarrollo con IA", 13,
     RGBColor(0xCF, 0xD6, 0xBE), space_after=10)
para(tf2, "Gonzalo Cantarero Galvez", 20, ACCENT, bold=True, space_after=0)

# ----------------------------------------------------------------------------
# SLIDE 2 — Índice
# ----------------------------------------------------------------------------
s = new_slide()
bg(s)
header(s, "Índice")
agenda = [
    ("1", "Problema y Solución", "Contexto, propuesta de valor y diferenciadores"),
    ("2", "Stack Tecnológico", ".NET 10, Supabase, Razor, Alpine, HTMX, Chart.js"),
    ("3", "Arquitectura", "Clean Architecture, CQRS, DDD, modelo de datos"),
    ("4", "Funcionalidades por Pantalla", "Recorrido completo tab por tab"),
    ("5", "Metodología de Desarrollo", "SDD, TDD y asistencia con IA"),
    ("6", "Testing y CI/CD", "Playwright E2E y pipeline GitHub Actions"),
    ("7", "Demo", "Aplicación en vivo con usuario de prueba"),
    ("8", "Conclusiones", "Logros, métricas y trabajo futuro"),
]
col_w = Inches(6.0)
x0, y0 = Inches(0.7), Inches(1.45)
for i, (n, t, d) in enumerate(agenda):
    col = i // 4
    row = i % 4
    x = x0 + col * Inches(6.2)
    y = y0 + row * Inches(1.32)
    chip = rrect(s, x, y, Inches(0.75), Inches(0.75), PRIMARY)
    ctf = chip.text_frame
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = ctf.paragraphs[0].add_run()
    set_run(r, n, 26, WHITE, bold=True)
    tb, tf = textbox(s, x + Inches(0.95), y - Inches(0.02), Inches(5.0),
                     Inches(0.85), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, 15, TEXT_DARK, bold=True, first=True, space_after=1)
    para(tf, d, 11, TEXT_MUTED, space_after=0)
footer(s, 2)

# ----------------------------------------------------------------------------
# Section divider helper
# ----------------------------------------------------------------------------
def divider(number, title, subtitle):
    s = new_slide()
    bg(s, PRIMARY)
    rect(s, 0, 0, Inches(0.35), SH, ACCENT)
    tb, tf = textbox(s, Inches(1.2), Inches(2.5), Inches(11), Inches(2.5),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, f"{number}", 90, RGBColor(0x7E, 0x8F, 0x54), bold=True, first=True,
         space_after=0)
    para(tf, title, 40, WHITE, bold=True, space_after=6)
    para(tf, subtitle, 16, RGBColor(0xE8, 0xE8, 0xDD), space_after=0)
    return s

# ----------------------------------------------------------------------------
# SECTION 1 — Problema y Solución
# ----------------------------------------------------------------------------
divider("01", "Problema y Solución", "¿Qué problema resuelve SauronSheet y para quién?")

# Slide: ¿Qué es SauronSheet?
s = new_slide(); bg(s); header(s, "¿Qué es SauronSheet?", 1)
tb, tf = textbox(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(1.5))
para(tf, "Aplicación web de finanzas personales que ayuda a tomar el control de "
         "gastos e ingresos de forma sencilla y visual.", 17, TEXT_DARK, first=True,
     space_after=4)
para(tf, "El nombre evoca «un ojo que todo lo ve» sobre tus finanzas — Sauron "
         "de El Señor de los Anillos + Sheet (hoja de cálculo).", 13, TEXT_MUTED,
     italic=True, space_after=0)
cards = [
    ("📥", "Importa", "Extractos bancarios desde Excel, sin introducir datos a mano"),
    ("📊", "Visualiza", "Dashboard interactivo con gráficos, tendencias y comparativas"),
    ("🎯", "Presupuesta", "Presupuestos por categoría con semáforo visual"),
    ("📈", "Analiza", "Análisis anual con score, anomalías y predicciones"),
]
cw = Inches(2.95); gap = Inches(0.15); x0 = Inches(0.7); y = Inches(3.2)
for i, (ic, t, d) in enumerate(cards):
    x = x0 + i * (cw + gap)
    card = rrect(s, x, y, cw, Inches(2.6), CARD_BG, line=CARD_BORDER)
    tb, tf = textbox(s, x + Inches(0.2), y + Inches(0.25), cw - Inches(0.4),
                     Inches(2.2), anchor=MSO_ANCHOR.TOP)
    para(tf, ic, 34, PRIMARY, first=True, align=PP_ALIGN.CENTER, space_after=6)
    para(tf, t, 17, PRIMARY, bold=True, align=PP_ALIGN.CENTER, space_after=6)
    para(tf, d, 12, TEXT_MUTED, align=PP_ALIGN.CENTER, space_after=0)
footer(s, 4)

# Slide: El problema
s = new_slide(); bg(s); header(s, "¿Por qué es difícil gestionar las finanzas?", 1)
problems = [
    ("Datos manuales", "Introducir cada movimiento a mano es tedioso y propenso a errores."),
    ("Sin visión global", "No hay una foto clara de ingresos frente a gastos."),
    ("Presupuestos olvidados", "Sin alertas visuales, los límites se ignoran."),
    ("Falta de contexto", "No hay análisis histórico, tendencias ni predicciones."),
]
y = Inches(1.55)
for i, (t, d) in enumerate(problems):
    yy = y + i * Inches(1.28)
    dot = rrect(s, Inches(0.75), yy, Inches(0.55), Inches(0.55), EXPENSE)
    dtf = dot.text_frame; dtf.paragraphs[0].alignment = PP_ALIGN.CENTER
    dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rr = dtf.paragraphs[0].add_run(); set_run(rr, "✕", 18, WHITE, bold=True)
    tb, tf = textbox(s, Inches(1.55), yy - Inches(0.05), Inches(11), Inches(1.1),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, 17, TEXT_DARK, bold=True, first=True, space_after=1)
    para(tf, d, 13, TEXT_MUTED, space_after=0)
footer(s, 5)

# Slide: La solución
s = new_slide(); bg(s); header(s, "SauronSheet — Un ojo que todo lo ve", 1)
sol = [
    "Importación automática de extractos bancarios (Excel ING)",
    "Dashboard interactivo con KPIs, gráficos y tendencias",
    "Presupuestos con semáforo verde / amarillo / rojo / sobrecoste",
    "Análisis anual completo con 7 secciones analíticas",
    "Score de salud financiera (0–100) con desglose",
    "Detección de anomalías y predicciones (regresión lineal)",
    "Multi-idioma (ES / EN) con cambio en vivo",
    "Multi-tenant: cada usuario ve solo sus datos (RLS)",
]
x0 = Inches(0.75); y0 = Inches(1.55); cw = Inches(5.85)
for i, txt in enumerate(sol):
    col = i // 4; row = i % 4
    x = x0 + col * Inches(6.1); yy = y0 + row * Inches(1.22)
    card = rrect(s, x, yy, cw, Inches(1.05), CARD_BG, line=CARD_BORDER)
    chk = rrect(s, x + Inches(0.18), yy + Inches(0.24), Inches(0.55), Inches(0.55), INCOME)
    ctf = chk.text_frame; ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rr = ctf.paragraphs[0].add_run(); set_run(rr, "✓", 18, WHITE, bold=True)
    tb, tf = textbox(s, x + Inches(0.95), yy, cw - Inches(1.1), Inches(1.05),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, txt, 13, TEXT_DARK, bold=True, first=True, space_after=0)
footer(s, 6)

# Slide: Para quién / diferenciadores
s = new_slide(); bg(s); header(s, "¿Para quién? ¿Qué nos diferencia?", 1)
bullet_card(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(4.9), "Usuarios objetivo",
    ["Personas que quieren entender dónde se va su dinero sin esfuerzo manual",
     "Usuarios de banca online que descargan extractos bancarios",
     "Quien quiera establecer presupuestos y recibir alertas visuales de sobrecoste"],
    icon="👥", body_size=14)
bullet_card(s, Inches(6.75), Inches(1.5), Inches(5.9), Inches(4.9), "Diferenciadores",
    ["Importación inteligente con resolución automática de categorías del banco",
     "Análisis anual completo: 7 secciones analíticas",
     "Multi-idioma con cambio en vivo (ES / EN)",
     "Multi-tenant aislado por Row-Level Security",
     "Presupuestos con semáforo y barras de progreso",
     "Trazabilidad total de importación por transacción"],
    icon="⭐", title_color=ACCENT, body_size=14)
footer(s, 7)

# ----------------------------------------------------------------------------
# SECTION 2 — Stack Tecnológico
# ----------------------------------------------------------------------------
divider("02", "Stack Tecnológico", "Tecnologías, librerías e infraestructura")

def tech_table(slide, x, y, w, rows, header_row, col_ws):
    n = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n, len(header_row), x, y, w,
                                       Inches(0.5 * n)).table
    for j, cw in enumerate(col_ws):
        tbl_shape.columns[j].width = cw
    # header
    for j, htext in enumerate(header_row):
        c = tbl_shape.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = PRIMARY
        c.margin_left = Inches(0.1); c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
        p = c.text_frame.paragraphs[0]; r = p.add_run()
        set_run(r, htext, 12, WHITE, bold=True)
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = tbl_shape.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF0, 0xEE, 0xE6)
            c.margin_left = Inches(0.1); c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            p = c.text_frame.paragraphs[0]; r = p.add_run()
            set_run(r, val, 11, TEXT_DARK, bold=(j == 0))
    return tbl_shape

# Backend
s = new_slide(); bg(s); header(s, "Stack — Backend", 2)
tech_table(s, Inches(0.7), Inches(1.5), Inches(11.9),
    [[".NET", "10", "Framework web y API — minimal APIs, Razor Pages"],
     ["C#", "13", "Lenguaje principal"],
     ["MediatR", "Última", "Patrón CQRS — commands y queries"],
     ["Sentry", "SDK .NET", "Observabilidad: trazabilidad, errores, métricas"],
     ["ExcelDataReader", "—", "Parseo de ficheros Excel (.xls / .xlsx)"]],
    ["Tecnología", "Versión", "Propósito"],
    [Inches(2.6), Inches(1.8), Inches(7.5)])
tb, tf = textbox(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.2))
para(tf, "Observabilidad 100% vía Sentry: pipeline único de tracing, errores, logs "
         "y métricas. Nunca se usa Console.WriteLine.", 13, TEXT_MUTED, italic=True,
     first=True)
footer(s, 9)

# Frontend
s = new_slide(); bg(s); header(s, "Stack — Frontend", 2)
tech_table(s, Inches(0.7), Inches(1.5), Inches(11.9),
    [["Razor Pages", ".NET 10", "Renderizado server-side, PageModel pattern"],
     ["MDBootstrap 5", "9.2", "UI Kit Material Design (CDN)"],
     ["Alpine.js", "3.14", "Reactividad declarativa en el DOM"],
     ["HTMX", "2.0", "AJAX desde atributos HTML"],
     ["Chart.js", "Última", "Gráficos interactivos"],
     ["Flatpickr", "Última", "Selector de fechas accesible"],
     ["Font Awesome", "6.5", "Iconografía funcional"]],
    ["Tecnología", "Versión", "Propósito"],
    [Inches(2.6), Inches(1.8), Inches(7.5)])
footer(s, 10)

# BD e Infra
s = new_slide(); bg(s); header(s, "Stack — Base de Datos e Infraestructura", 2)
bullet_card(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(4.9),
    "Base de datos y Auth",
    ["Supabase (PostgreSQL): base de datos relacional gestionada",
     "Supabase Auth: registro, login, JWT, refresh tokens",
     "Row-Level Security: aislamiento multi-tenant en BD",
     "Supabase CLI: migraciones y gestión del proyecto",
     "20 migraciones SQL versionadas"],
    icon="🗄️", body_size=14)
bullet_card(s, Inches(6.75), Inches(1.5), Inches(5.9), Inches(4.9),
    "Infraestructura",
    ["Azure App Service: hosting de la aplicación .NET",
     "Supabase Cloud: base de datos y autenticación",
     "Sentry: observabilidad unificada (pipeline único)",
     "GitHub Actions: CI/CD (tests, E2E, migraciones, deploy)",
     "GitHub: repositorio y control de versiones"],
    icon="☁️", title_color=ACCENT, body_size=14)
footer(s, 11)

# ----------------------------------------------------------------------------
# SECTION 3 — Arquitectura
# ----------------------------------------------------------------------------
divider("03", "Arquitectura", "Clean Architecture, CQRS, DDD y modelo de datos")

# Clean Architecture
s = new_slide(); bg(s); header(s, "Clean Architecture + CQRS", 3)
layers = [
    ("FRONTEND (Razor Pages)", "Pages · ViewModels · wwwroot — depende de Application", PRIMARY_LT),
    ("APPLICATION (CQRS)", "Commands · Queries · DTOs · Behaviors — depende de Domain", PRIMARY),
    ("DOMAIN (Core)", "Entities · ValueObjects · Services · Specifications — cero dependencias", PRIMARY_DARK),
    ("INFRASTRUCTURE", "Supabase · Auth · Excel · Sentry — implementa interfaces del Domain", PRIMARY_LT),
]
y = Inches(1.5)
for i, (t, d, col) in enumerate(layers):
    yy = y + i * Inches(1.12)
    band = rrect(s, Inches(1.2), yy, Inches(9.2), Inches(0.95), col)
    tb, tf = textbox(s, Inches(1.45), yy, Inches(8.7), Inches(0.95),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, 15, WHITE, bold=True, first=True, space_after=1)
    para(tf, d, 11, RGBColor(0xEC, 0xEF, 0xE2), space_after=0)
# arrow indicating inward dependency
tb, tf = textbox(s, Inches(10.6), Inches(1.6), Inches(2.5), Inches(4.5),
                 anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Las dependencias\napuntan\nHACIA ADENTRO", 13, PRIMARY, bold=True,
     first=True, align=PP_ALIGN.CENTER, space_after=8)
para(tf, "El Domain no conoce\nninguna otra capa", 11, TEXT_MUTED,
     align=PP_ALIGN.CENTER, space_after=0)
footer(s, 13)

# CQRS
s = new_slide(); bg(s); header(s, "Patrón CQRS con MediatR", 3)
bullet_card(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(3.4),
    "Commands (escritura)",
    ["CreateTransactionCommand → Transaction",
     "ImportTransactionsCommand → Transaction[]",
     "CreateBudgetCommand → Budget",
     "UpdateTransactionCommand / DeleteTransactionCommand"],
    icon="✍️", body_size=13)
bullet_card(s, Inches(6.75), Inches(1.5), Inches(5.9), Inches(3.4),
    "Queries (lectura)",
    ["GetTransactionsQuery → TransactionDto[]",
     "GetSpendingByCategoryQuery → CategorySpending[]",
     "GetAnnualDashboardQuery → AnnualDashboardResult",
     "GetBudgetMetricsQuery → BudgetMetricsDto[]"],
    icon="🔍", title_color=ACCENT, body_size=13)
bullet_card(s, Inches(0.7), Inches(5.1), Inches(11.95), Inches(1.35),
    "Pipeline behaviors (cross-cutting)",
    ["TenantScopingBehavior — inyecta el UserId en cada request automáticamente        "
     "•  SentryTracingBehavior — traza cada comando / query en Sentry"],
    icon="⚙️", body_size=13)
footer(s, 14)

# DDD
s = new_slide(); bg(s); header(s, "Domain-Driven Design", 3)
bullet_card(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(2.55),
    "Aggregate Roots",
    ["Transaction — Amount (Money), Date, CategoryId, Balance",
     "Category — Name, Type, Source, IsAutoCreated",
     "Budget — Limit (Money), Period, Status, fechas",
     "Subcategory · ImportBatch"],
    icon="🧩", body_size=12)
bullet_card(s, Inches(6.75), Inches(1.5), Inches(5.9), Inches(2.55),
    "Value Objects inmutables",
    ["Money (decimal, Currency) — validación + aritmética",
     "DateRange · BudgetPeriod (Year, Month)",
     "TransactionId · UserId · CategoryName (tipados fuertes)",
     "BudgetStatusLevel — Verde / Amarillo / Rojo / Sobrecoste"],
    icon="💎", title_color=ACCENT, body_size=12)
bullet_card(s, Inches(0.7), Inches(4.2), Inches(5.9), Inches(2.25),
    "Domain Services",
    ["CategoryService — nombres únicos, categorías del sistema",
     "BudgetCalculationService — % usado, nivel de estado",
     "BudgetService — reglas de negocio de presupuestos"],
    icon="🛠️", body_size=12)
bullet_card(s, Inches(6.75), Inches(4.2), Inches(5.9), Inches(2.25),
    "Specifications",
    ["Por rango de fechas / categoría / importe",
     "Por palabra clave en descripción / usuario",
     "CompositeSpecification — combinación AND / OR"],
    icon="📐", title_color=ACCENT, body_size=12)
footer(s, 15)

# Modelo de datos
s = new_slide(); bg(s); header(s, "Modelo de Datos en Supabase (PostgreSQL)", 3)
tables = [
    ("user_profiles", "Perfiles de usuario extendidos"),
    ("categories", "Tipo (Income/Expense), source (System/User), únicas por nombre"),
    ("subcategories", "Clasificación detallada dentro de una categoría"),
    ("transactions", "Fecha, descripción, amount signado, balance, trazabilidad"),
    ("budgets", "limit_amount, period (YYYY-MM), status, fechas de vigencia"),
    ("import_batches", "file_name, status, transaction_count, errores (JSON)"),
    ("bank_category_translations", "Traducción banco → categoría SauronSheet"),
]
y = Inches(1.5)
for i, (t, d) in enumerate(tables):
    yy = y + i * Inches(0.66)
    chip = rrect(s, Inches(0.7), yy, Inches(3.6), Inches(0.55), PRIMARY)
    ctf = chip.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = Inches(0.15)
    r = ctf.paragraphs[0].add_run(); set_run(r, t, 12, WHITE, bold=True)
    tb, tf = textbox(s, Inches(4.5), yy, Inches(8.2), Inches(0.55),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, d, 12, TEXT_DARK, first=True)
tb, tf = textbox(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.7))
para(tf, "RLS habilitado en todas las tablas: filtran por user_id = auth.uid(). "
         "Cliente Supabase inicializado por request con el JWT → aislamiento total.",
     12, TEXT_MUTED, italic=True, first=True)
footer(s, 16)

# ----------------------------------------------------------------------------
# SECTION 4 — Funcionalidades por Pantalla (TABS)
# ----------------------------------------------------------------------------
divider("04", "Funcionalidades por Pantalla", "Recorrido completo por cada tab de la aplicación")

def feature_slide(page, icon, title, screens, features, note=None):
    s = new_slide(); bg(s); header(s, f"{icon}  {title}", 4)
    # screens chip row
    tb, tf = textbox(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.5))
    r = tf.paragraphs[0].add_run()
    set_run(r, "Pantallas:  ", 12, PRIMARY, bold=True)
    r2 = tf.paragraphs[0].add_run()
    set_run(r2, screens, 12, TEXT_MUTED, italic=True)
    # features list in card
    card = rrect(s, Inches(0.7), Inches(1.95), Inches(11.95),
                 Inches(4.35) if note else Inches(4.7), CARD_BG, line=CARD_BORDER)
    tb, tf = textbox(s, Inches(1.0), Inches(2.2), Inches(11.4),
                     Inches(3.9) if note else Inches(4.3))
    first = True
    for f in features:
        p = para(tf, "•  " + f, 14, TEXT_DARK, first=first, space_after=8)
        first = False
    if note:
        tb, tf = textbox(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.55))
        para(tf, note, 12, PRIMARY, italic=True, bold=True, first=True)
    footer(s, page)
    return s

# 4.1 Autenticación
feature_slide(18, "👤", "Autenticación Multi-Usuario",
    "Auth/Login · Auth/Register · Auth/Logout",
    ["Registro e inicio de sesión con email y contraseña (Supabase Auth)",
     "JWT con refresh tokens almacenados en cookies seguras (HttpOnly, SameSite Strict)",
     "Aislamiento completo de datos por usuario mediante RLS en Supabase",
     "Cierre de sesión con limpieza de cookies"],
    note="Cada usuario ve exclusivamente sus propios datos — arquitectura multi-tenant.")

# 4.2 Dashboard
feature_slide(19, "📊", "Dashboard Principal",
    "Dashboard",
    ["KPIs animados: ingresos totales, gastos totales, neto y nº de transacciones",
     "Gráfico de gastos por categoría (barras apiladas por mes)",
     "Gráfico de tendencias mensuales (ingresos vs gastos)",
     "Comparativa año contra año y transacciones recientes",
     "Widget de estado de presupuestos con semáforo y barras de progreso",
     "Filtro por periodo: Todo · Este mes · Últimos 3 meses · Este año"],
    note="Punto de entrada visual — foto completa de la salud financiera del usuario.")

# 4.3 Transacciones
feature_slide(20, "🔍", "Gestión de Transacciones",
    "Transactions/Index · Add · Edit · Search",
    ["Listado paginado con filtros por fecha, categoría, importe y origen",
     "Búsqueda por palabra clave en la descripción",
     "Alta manual de transacciones",
     "Edición de descripción, importe, fecha y categoría",
     "Eliminación individual o masiva",
     "Vista detalle con toda la información de importación (trazabilidad)"],
    note="CRUD completo sobre el núcleo de datos de la aplicación.")

# 4.4 Importación
feature_slide(21, "📄", "Importación de Extractos Bancarios",
    "Transactions/Upload",
    ["Subida de ficheros Excel (.xls / .xlsx) con extractos bancarios",
     "Parseo automático de movimientos: fecha, descripción, importe y saldo",
     "Detección de duplicados por saldo y fecha para evitar repeticiones",
     "Resolución automática de categorías basada en la subcategoría del banco",
     "Progreso de la importación en tiempo real",
     "Trazabilidad: origen (banco), fichero, fecha de importación y método de categorización"],
    note="Funcionalidad estrella: elimina la introducción manual de datos.")

# 4.5 Categorías
feature_slide(22, "🏷️", "Categorías y Subcategorías",
    "Categories/Index · Categories/Subcategories",
    ["Categorías por defecto del sistema (Comida, Transporte, Servicios, Otros)",
     "Categorías personalizadas por usuario",
     "Subcategorías para clasificación detallada",
     "Categorización manual de transacciones no clasificadas",
     "Resolución automática desde la subcategoría del banco"],
    note="Las categorías del sistema son inmutables (IsSystemDefault) y sembradas por defecto.")

# 4.6 Presupuestos
feature_slide(23, "💰", "Gestión de Presupuestos",
    "Budgets/Index · Create · Edit · Metrics · History · Comparison",
    ["Creación de presupuestos mensuales por categoría con fecha de inicio/fin",
     "Edición de límite, periodo y fechas de vigencia",
     "Semáforo visual: Verde (≤75%) · Amarillo (75-90%) · Rojo (90-100%) · Sobrecoste (>100%)",
     "Barras de progreso con importe gastado vs límite",
     "Histórico de presupuestos anteriores",
     "Métricas: gastado, límite acumulado, porcentaje usado y restante"],
    note="El semáforo cambia de color automáticamente según el porcentaje consumido.")

# 4.7 Análisis Anual (parte 1)
s = new_slide(); bg(s); header(s, "📈  Análisis Anual — Secciones 1 a 4", 4)
tb, tf = textbox(s, Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.45))
r = tf.paragraphs[0].add_run(); set_run(r, "Pantalla:  ", 12, PRIMARY, bold=True)
r2 = tf.paragraphs[0].add_run(); set_run(r2, "Analysis/Annual", 12, TEXT_MUTED, italic=True)
secs1 = [
    ("1 · Resumen Ejecutivo", "Ingresos, gastos, neto, ahorro y tasa de ahorro del año. Variación vs año anterior y ranking."),
    ("2 · Salud Financiera", "Score global 0–100 con anillo visual. Sub-scores: ahorro, estabilidad, dependencia, balance, tendencia."),
    ("3 · Ratios Financieros", "Tasa de ahorro, ingreso/gasto mensual medio, nº transacciones y % de costes fijos."),
    ("4 · Comparativa Interanual (YoY)", "Año actual vs anterior lado a lado: variación absoluta y porcentual de todas las métricas."),
]
y = Inches(1.9)
for i, (t, d) in enumerate(secs1):
    yy = y + i * Inches(1.18)
    card = rrect(s, Inches(0.7), yy, Inches(11.95), Inches(1.02), CARD_BG, line=CARD_BORDER)
    num = rrect(s, Inches(0.7), yy, Inches(0.18), Inches(1.02), ACCENT)
    tb, tf = textbox(s, Inches(1.05), yy, Inches(11.4), Inches(1.02),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, 15, PRIMARY, bold=True, first=True, space_after=1)
    para(tf, d, 12, TEXT_DARK, space_after=0)
footer(s, 24)

# 4.7 Análisis Anual (parte 2)
s = new_slide(); bg(s); header(s, "📈  Análisis Anual — Secciones 5 a 7", 4)
secs2 = [
    ("5 · Tendencia y Distribución Mensual", "Gráfico de tendencia mensual (ingresos vs gastos), distribución fijo/variable y mejor mes."),
    ("6 · Categorías y Comparativa", "Desglose de gastos por categoría con gráfico donut, tabla de tendencias y comparativa YoY."),
    ("7 · Anomalías y Predicciones", "Detección estadística (Z-score > 3) y predicciones deterministas por regresión lineal (R² como confianza)."),
]
y = Inches(1.55)
for i, (t, d) in enumerate(secs2):
    yy = y + i * Inches(1.35)
    card = rrect(s, Inches(0.7), yy, Inches(11.95), Inches(1.15), CARD_BG, line=CARD_BORDER)
    rrect(s, Inches(0.7), yy, Inches(0.18), Inches(1.15), ACCENT)
    tb, tf = textbox(s, Inches(1.05), yy, Inches(11.4), Inches(1.15),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, 16, PRIMARY, bold=True, first=True, space_after=2)
    para(tf, d, 13, TEXT_DARK, space_after=0)
tb, tf = textbox(s, Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.7))
para(tf, "7 secciones analíticas que convierten datos crudos en decisiones financieras "
         "accionables, con predicciones a futuro.", 12, TEXT_MUTED, italic=True, first=True)
footer(s, 25)

# 4.8 i18n + diseño
s = new_slide(); bg(s); header(s, "🌐  Internacionalización y Diseño", 4)
bullet_card(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(4.9),
    "Internacionalización (i18n)",
    ["Interfaz completa en español (es-ES) e inglés (en-US)",
     "Selector de idioma en el menú de navegación",
     "Persistencia de la elección mediante cookie",
     "Localización de Chart.js, Flatpickr, fechas y moneda",
     "Sistema de recursos .resx con SharedResources"],
    icon="🌐", body_size=14)
bullet_card(s, Inches(6.75), Inches(1.5), Inches(5.9), Inches(4.9),
    "Diseño Olive + Responsivo",
    ["Color primario Olive Green #556B2F (estabilidad y crecimiento)",
     "Canvas #f8f9fa con tarjetas blancas elevadas",
     "Semántica: verde ingresos, rojo gastos",
     "Navegación adaptable: menú colapsable y offcanvas en móvil",
     "Cuadrículas 1 columna (móvil) → 4 columnas (escritorio)",
     "Alpine.js + HTMX + Chart.js para interactividad"],
    icon="🎨", title_color=ACCENT, body_size=13)
footer(s, 26)

# ----------------------------------------------------------------------------
# SECTION 5 — Metodología
# ----------------------------------------------------------------------------
divider("05", "Metodología de Desarrollo", "SDD, TDD y desarrollo asistido por IA")

# SDD
s = new_slide(); bg(s); header(s, "Spec-Driven Development (SDD)", 5)
phases = ["Proposal", "Spec", "Design", "Tasks", "Apply", "Verify", "Archive"]
x = Inches(0.55); y = Inches(1.7); pw = Inches(1.62); ph = Inches(0.95)
for i, ph_name in enumerate(phases):
    xx = x + i * Inches(1.75)
    box = rrect(s, xx, y, pw, ph, PRIMARY if i % 2 == 0 else PRIMARY_LT)
    btf = box.text_frame; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = btf.paragraphs[0].add_run(); set_run(r, ph_name, 12, WHITE, bold=True)
    if i < len(phases) - 1:
        tb, tf = textbox(s, xx + pw - Inches(0.02), y, Inches(0.2), ph,
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, "›", 20, ACCENT, bold=True, first=True, align=PP_ALIGN.CENTER)
tb, tf = textbox(s, Inches(0.7), Inches(3.1), Inches(11.9), Inches(3.2))
steps = [
    "Exploración — investigación del código y requisitos",
    "Propuesta — alcance, enfoque y tradeoffs",
    "Especificación — requisitos detallados y escenarios de aceptación",
    "Diseño técnico — arquitectura, patrones y decisiones",
    "Tareas — desglose en unidades atómicas implementables",
    "Implementación — codificación con TDD estricto",
    "Verificación — tests, revisión y validación contra la spec",
    "Archivo — cierre de fase y registro de decisiones",
]
first = True
for st in steps:
    para(tf, "•  " + st, 13, TEXT_DARK, first=first, space_after=6); first = False
footer(s, 28)

# TDD + pyramid
s = new_slide(); bg(s); header(s, "TDD y Pirámide de Testing", 5)
bullet_card(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(2.5),
    "TDD (Test-Driven Development)",
    ["Flujo: Test Rojo → Implementación → Test Verde → Refactor",
     "Cobertura mínima: 80% Domain · 70% Application",
     "Los tests preceden a la implementación"],
    icon="🧪", body_size=14)
# pyramid on the right
px = Inches(7.4); py = Inches(1.6)
pyr = [
    ("E2E — Playwright (browser automation)", Inches(3.0), ACCENT),
    ("Integration — xUnit + Moq + in-memory", Inches(4.2), PRIMARY_LT),
    ("Unit (Domain) — xUnit", Inches(5.4), PRIMARY),
]
for i, (label, w, col) in enumerate(pyr):
    yy = py + i * Inches(0.85)
    xx = px + (Inches(5.4) - w) / 2
    band = rrect(s, xx, yy, w, Inches(0.72), col)
    btf = band.text_frame; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = btf.paragraphs[0].add_run(); set_run(r, label, 10, WHITE, bold=True)
bullet_card(s, Inches(0.7), Inches(4.2), Inches(5.9), Inches(2.25),
    "Resultados",
    ["910 tests (unidad, integración, frontend, infraestructura)",
     "10 specs E2E con Playwright (~60 escenarios)",
     "Ejecución en CI antes de cada deploy"],
    icon="✅", title_color=ACCENT, body_size=14)
tb, tf = textbox(s, Inches(7.4), Inches(4.4), Inches(5.2), Inches(2.0),
                 anchor=MSO_ANCHOR.MIDDLE)
para(tf, "910", 60, PRIMARY, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=0)
para(tf, "tests automatizados", 15, TEXT_MUTED, align=PP_ALIGN.CENTER, space_after=0)
footer(s, 29)

# IA
s = new_slide(); bg(s); header(s, "Desarrollo Asistido por IA — Gentle AI", 5)
tb, tf = textbox(s, Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.9))
para(tf, "Todo el desarrollo se ha realizado con Gentle AI sobre opencode, "
         "orquestando el flujo SDD con un ecosistema de herramientas de "
         "optimización de tokens y memoria persistente.", 14, TEXT_DARK, first=True)
tools = [
    ("Headroom", "Compresión de contexto en ventana — hasta 90% menos tokens en respuestas largas y logs"),
    ("Serena", "Análisis semántico con LSP — navegación precisa de símbolos sin depender de grep"),
    ("RTK (Rust Token Killer)", "Compresión de comandos de terminal — 60–99% menos tokens en git, dotnet, grep"),
    ("Engram", "Memoria persistente entre sesiones — decisiones, bugs y configuraciones"),
]
x0 = Inches(0.7); y0 = Inches(2.5); cw = Inches(5.9); ch = Inches(1.75)
for i, (t, d) in enumerate(tools):
    col = i % 2; row = i // 2
    x = x0 + col * Inches(6.05); yy = y0 + row * Inches(1.95)
    card = rrect(s, x, yy, cw, ch, CARD_BG, line=CARD_BORDER)
    tb, tf = textbox(s, x + Inches(0.25), yy + Inches(0.18), cw - Inches(0.5), ch - Inches(0.36))
    para(tf, t, 16, PRIMARY, bold=True, first=True, space_after=4)
    para(tf, d, 12, TEXT_MUTED, space_after=0)
footer(s, 30)

# ----------------------------------------------------------------------------
# SECTION 6 — Testing y CI/CD
# ----------------------------------------------------------------------------
divider("06", "Testing y CI/CD", "Playwright E2E y pipeline de despliegue continuo")

# Playwright
s = new_slide(); bg(s); header(s, "Tests E2E con Playwright", 6)
bullet_card(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(4.9),
    "¿Por qué Playwright?",
    ["Multi-navegador: Chromium, Firefox, WebKit con una API",
     "Auto-wait inteligente — sin sleep() arbitrarios",
     "Selectores data-testid — inmunes al cambio de idioma",
     "Trace viewer, screenshots y vídeo en fallo",
     "Integración nativa con GitHub Actions"],
    icon="🎭", body_size=14)
bullet_card(s, Inches(6.75), Inches(1.5), Inches(5.9), Inches(4.9),
    "Flujos cubiertos (9 specs, ~60 tests)",
    ["Cambio de idioma ES↔EN y persistencia",
     "Login válido/inválido, registro, logout",
     "Importación de extracto ING con duplicados",
     "CRUD de presupuestos y semáforo",
     "Edición de transacciones",
     "Ciclo de vida de categorías y subcategorías",
     "Navegación del análisis anual (7 secciones)"],
    icon="📋", title_color=ACCENT, body_size=13)
footer(s, 32)

# CI/CD
s = new_slide(); bg(s); header(s, "Pipeline CI/CD con GitHub Actions", 6)
steps = [
    ("1", "Build & Test", "dotnet restore → build → test (910 tests)"),
    ("2", "E2E Tests", "Playwright Chromium — 9 specs con credenciales de test"),
    ("3", "Migraciones", "supabase link → db push (antes del deploy)"),
    ("4", "Publish", "dotnet publish con limpieza de .map y .pdb"),
    ("5", "Deploy a Azure", "Login OIDC → azure/webapps-deploy a App Service"),
]
y = Inches(1.55)
for i, (n, t, d) in enumerate(steps):
    yy = y + i * Inches(0.92)
    chip = rrect(s, Inches(0.7), yy, Inches(0.7), Inches(0.7), PRIMARY)
    ctf = chip.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = ctf.paragraphs[0].add_run(); set_run(r, n, 20, WHITE, bold=True)
    card = rrect(s, Inches(1.6), yy, Inches(11.05), Inches(0.7), CARD_BG, line=CARD_BORDER)
    tb, tf = textbox(s, Inches(1.85), yy, Inches(10.6), Inches(0.7),
                     anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); set_run(r, t + " — ", 14, PRIMARY, bold=True)
    r2 = p.add_run(); set_run(r2, d, 13, TEXT_DARK)
tb, tf = textbox(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.7))
para(tf, "Trigger: push a ramas RELEASE/Sauron.Sheet/*.  Rollback: revertir el commit "
         "y pushear → redeploy automático de la versión anterior.", 12, TEXT_MUTED,
     italic=True, first=True)
footer(s, 33)

# ----------------------------------------------------------------------------
# SECTION 7 — Demo
# ----------------------------------------------------------------------------
divider("07", "Demo", "La aplicación en vivo con usuario de prueba")

s = new_slide(); bg(s); header(s, "Demo en Vivo", 7)
card = rrect(s, Inches(2.0), Inches(1.7), Inches(9.3), Inches(3.9), CARD_BG, line=CARD_BORDER)
tb, tf = textbox(s, Inches(2.4), Inches(2.0), Inches(8.5), Inches(3.4))
para(tf, "🌐  Aplicación desplegada", 16, PRIMARY, bold=True, first=True, space_after=3)
para(tf, "sauronsheet-akd4gkewdpbtbgea.spaincentral-01.azurewebsites.net", 13,
     TEXT_DARK, space_after=16)
para(tf, "👤  Usuario de prueba", 16, PRIMARY, bold=True, space_after=3)
para(tf, "Email:  demo@sauronsheet.app", 14, TEXT_DARK, space_after=2)
para(tf, "Contraseña:  Demo1234!", 14, TEXT_DARK, space_after=16)
para(tf, "El usuario demo incluye datos de ejemplo: transacciones importadas, "
         "presupuestos activos y análisis anual con datos históricos (2024–2026).",
     12, TEXT_MUTED, italic=True, space_after=0)
footer(s, 35)

# ----------------------------------------------------------------------------
# SECTION 8 — Conclusiones
# ----------------------------------------------------------------------------
divider("08", "Conclusiones", "Logros, métricas y trabajo futuro")

# Logros + métricas
s = new_slide(); bg(s); header(s, "Logros y Métricas del Proyecto", 8)
metrics = [("910", "tests"), ("29", "migraciones/pantallas"), ("4", "capas"), ("2", "idiomas")]
mx = Inches(0.7); my = Inches(1.5); mw = Inches(2.9)
metric_vals = [("910", "tests automatizados"), ("10", "specs E2E"),
               ("7", "secciones de análisis"), ("2", "idiomas ES / EN")]
for i, (v, l) in enumerate(metric_vals):
    x = mx + i * Inches(3.05)
    card = rrect(s, x, my, mw, Inches(1.5), PRIMARY if i % 2 == 0 else PRIMARY_LT)
    tb, tf = textbox(s, x, my + Inches(0.12), mw, Inches(1.3), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, v, 40, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=0)
    para(tf, l, 12, RGBColor(0xEC, 0xEF, 0xE2), align=PP_ALIGN.CENTER, space_after=0)
bullets_block(s, Inches(0.75), Inches(3.4), Inches(11.9), Inches(3.2), [
    "Aplicación web completa y desplegada en producción (Azure)",
    "Clean Architecture + CQRS + DDD aplicados correctamente",
    "TDD estricto en dominio con cobertura >80%",
    "Importación real de extractos bancarios con resolución de categorías",
    "Análisis anual avanzado con anomalías y predicciones",
    "Multi-idioma, multi-tenant y CI/CD completo con rollback",
], size=14, gap=7)
footer(s, 37)

# Roadmap
s = new_slide(); bg(s); header(s, "Trabajo Futuro", 8)
future = [
    ("📄", "Exportación a PDF / CSV de informes", "Media"),
    ("🔔", "Alertas push de presupuestos", "Media"),
    ("🏦", "Importación multi-banco (CaixaBank, Santander, BBVA)", "Media"),
    ("🌙", "Modo oscuro", "Baja"),
    ("🤖", "Categorización automática con ML (basada en histórico)", "Baja"),
    ("👥", "Presupuestos compartidos (grupo / pareja)", "Baja"),
]
y = Inches(1.55)
for i, (ic, t, pri) in enumerate(future):
    yy = y + i * Inches(0.85)
    card = rrect(s, Inches(0.7), yy, Inches(11.95), Inches(0.72), CARD_BG, line=CARD_BORDER)
    tb, tf = textbox(s, Inches(1.0), yy, Inches(9.5), Inches(0.72), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); set_run(r, ic + "   ", 16, PRIMARY)
    r2 = p.add_run(); set_run(r2, t, 14, TEXT_DARK, bold=True)
    prc = ACCENT if pri == "Media" else TEXT_MUTED
    pchip = rrect(s, Inches(11.0), yy + Inches(0.16), Inches(1.4), Inches(0.4), prc)
    ptf = pchip.text_frame; ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ptf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = ptf.paragraphs[0].add_run(); set_run(r, pri, 11, WHITE, bold=True)
footer(s, 38)

# ----------------------------------------------------------------------------
# SLIDE FINAL — Gracias
# ----------------------------------------------------------------------------
s = new_slide()
bg(s, PRIMARY_DARK)
rect(s, 0, 0, Inches(0.35), SH, ACCENT)
s.shapes.add_picture(LOGO, Inches(5.22), Inches(0.75), height=Inches(2.0))
tb, tf = textbox(s, Inches(1), Inches(3.0), Inches(11.33), Inches(1.4),
                 anchor=MSO_ANCHOR.MIDDLE)
para(tf, "¡Gracias!", 46, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=2)
para(tf, "Un ojo que todo lo ve sobre tus finanzas", 16, ACCENT, italic=True,
     align=PP_ALIGN.CENTER, space_after=0)
# links box
tb, tf = textbox(s, Inches(2.5), Inches(4.5), Inches(8.33), Inches(2.2),
                 anchor=MSO_ANCHOR.TOP)
para(tf, "Demo:  sauronsheet-akd4gkewdpbtbgea.spaincentral-01.azurewebsites.net",
     13, RGBColor(0xE8, 0xE8, 0xDD), first=True, align=PP_ALIGN.CENTER, space_after=5)
para(tf, "Usuario:  demo@sauronsheet.app  /  Demo1234!", 13,
     RGBColor(0xE8, 0xE8, 0xDD), align=PP_ALIGN.CENTER, space_after=14)
para(tf, "Gonzalo Cantarero Galvez", 20, ACCENT, bold=True,
     align=PP_ALIGN.CENTER, space_after=2)
para(tf, "Trabajo de Fin de Máster — Máster de Desarrollo con IA",
     12, RGBColor(0xCF, 0xD6, 0xBE), align=PP_ALIGN.CENTER, space_after=0)

# ----------------------------------------------------------------------------
prs.save(OUT)
print(f"OK - {len(prs.slides._sldIdLst)} slides")
print("Saved:", OUT)
