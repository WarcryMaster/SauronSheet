import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

p = Presentation(r'resources/SauronSheet_TFM_Presentation.pptx')
for i, slide in enumerate(p.slides, 1):
    texts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip().replace('\n', ' | ')
            if t:
                texts.append(t)
    title = texts[0][:75] if texts else '(no text)'
    body = ' /// '.join(texts[1:])[:120] if len(texts) > 1 else ''
    print(f"[{i:02d}] {title}")
    if body:
        print(f"      {body}")
